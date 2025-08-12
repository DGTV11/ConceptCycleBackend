import base64
from io import BytesIO

import fitz
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation

from llm import call_vlm


def vlm_process_image(b64_image, img_type):
    return call_vlm(
        """
# MISSION
Compress one image of school material (notes / slides / worksheet) into a token-efficient V-SPR that preserves title, ordered sections, key facts, equations (LaTeX), diagram relations, table sketch, and layout cues so another model can expand into detailed text notes.

# THEORY
Visual school artifacts combine dense text, structure (headings, lists), equations, and schematic diagrams. A small set of canonical, confidence-scored assertions (headlines, ordered section entries, short triples for diagrams/tables, LaTeX for math) activates a downstream model to reliably reconstruct full notes while minimizing tokens.

# METHODOLOGY
1. Detect page type and read structural order: Title → Section headings → subpoints → equations → diagrams → tables → captions/annotations.
2. Emit a compact ordered list of elements; keep original wording if OCR confidence ≥0.85, else paraphrase and append `CONF:0.??`.
3. Use these canonical element lines only (around 12–16 lines):
   `TITLE: <text>. CONF:x.xx`
   `SECTION n: <heading> | BULLETS: [short phrases] | CONF:x.xx`
   `EQUATION: $...$ | LOC: inline/center | CONF:x.xx`
   `DIAGRAM: <EntityA> -> <EntityB> (label). CONF:x.xx`
   `TABLE: header1,header2 | SAMPLE_ROWS: [r1,r2] | MORE_ROWS:N | CONF:x.xx`
4. For lists, compress to 2–4 distilled bullet statements preserving order and core meaning (use punctuation to separate clauses).
5. Preserve layout cues that change meaning: underlines, bold, indentation, numberings, arrows; encode as `EMPH: underline/bold/indent` with confidence.
6. Mark uncertain tokens inline as `[?word]` and give numeric confidence per line (0.00–1.00). End output with `OVERALL_CONF: x.xx`.
7. Strictly output only the compressed V-SPR (no freeform prose).
""",  # *https://github.com/daveshap/SparsePrimingRepresentations
        b64_image,
        img_type,
    )


def process_file(
    file: bytes, filename: str, content_type: str
):  # txt/md, images, pptx, word, pdf
    match content_type:
        case "txt":  # *applies for markdown obviously since its just txt
            # To read a FastAPI SpooledTemporaryFile (which is the underlying file object of an UploadFile) as text, the recommended approach is to use io.TextIOWrapper for proper encoding handling.
            return file.decode("utf-8")
        case "png":
            return vlm_process_image(base64.b64encode(file).decode("utf-8"), "png")
        case "jpeg":
            return vlm_process_image(base64.b64encode(file).decode("utf-8"), "jpeg")
        case "pptx":
            slides = Presentation(BytesIO(file)).slides

            slides_notes = ""
            for slide in slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_notes += shape.text + "\n"
                    if hasattr(shape, "image"):
                        slides_notes += (
                            "\n"
                            + "===IMAGE START==="
                            + vlm_process_image(
                                base64.b64encode(shape.image.blob).decode("utf-8"),
                                shape.image.content_type.replace("image/", ""),
                            )
                            + "===IMAGE END==="
                            + "\n"
                        )

            return slides_notes
        case "docx":
            doc = Document(BytesIO(file))
            docx_notes = ""

            for para in doc.paragraphs:
                for run in para.runs:
                    if run.text:
                        docx_notes += run.text

                    drawing_elems = run._element.findall(
                        ".//w:drawing", namespaces=run._element.nsmap
                    )
                    for drawing in drawing_elems:
                        namespaces = drawing.nsmap.copy() if drawing.nsmap else {}
                        if "a" not in namespaces:
                            namespaces["a"] = (
                                "http://schemas.openxmlformats.org/drawingml/2006/main"
                            )

                        blip = drawing.find(".//a:blip", namespaces=namespaces)
                        if blip is not None:
                            embed_rid = blip.get(qn("r:embed"))
                            image_part = doc.part.related_parts[embed_rid]

                            blob = image_part.blob
                            b64_blob = base64.b64encode(blob).decode("utf-8")
                            content_type = image_part.content_type  # e.g. image/png

                            docx_notes += (
                                "\n===IMAGE START===\n"
                                + vlm_process_image(
                                    b64_blob, content_type.replace("image/", "")
                                )
                                + "\n===IMAGE END===\n"
                            )

                docx_notes += "\n"

            return docx_notes
        case "pdf":
            doc = fitz.open("pdf", file)
            pdf_notes = ""
            for page in doc:
                pdf_notes += page.get_text("text") + "\n"

                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    blob = base_image["image"]
                    b64_blob = base64.b64encode(blob).decode("utf-8")
                    ext = base_image["ext"]
                    pdf_notes += (
                        "\n===IMAGE START===\n"
                        + vlm_process_image(b64_blob, ext)
                        + "\n===IMAGE END===\n"
                    )
            return pdf_notes
        case _:
            raise ValueError("Invalid content_type")
